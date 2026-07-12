"""Connector-line geometry extraction + native PPT line emission tests."""
from __future__ import annotations

import sys
import unittest
from pathlib import Path

import cv2
import numpy as np

REPO_ROOT = Path(__file__).resolve().parents[1]
PAGE_DIR = REPO_ROOT / "scripts" / "page"
SCRIPTS_DIR = REPO_ROOT / "scripts"
DECK_DIR = REPO_ROOT / "scripts" / "deck"
for path in (SCRIPTS_DIR, PAGE_DIR, DECK_DIR):
    if str(path) not in sys.path:
        sys.path.insert(0, str(path))

from layout.connector_geometry import extract_connector_geometry  # noqa: E402
from layout.icon_alpha import _line_art_matte  # noqa: E402


def _geometry_for(crop: np.ndarray):
    _fg, alpha = _line_art_matte(crop)
    return extract_connector_geometry(crop, alpha)


class ConnectorGeometryTests(unittest.TestCase):
    def test_horizontal_line(self) -> None:
        crop = np.full((24, 200, 3), 255, np.uint8)
        cv2.line(crop, (8, 12), (192, 12), (80, 80, 80), 3, cv2.LINE_AA)
        geom = _geometry_for(crop)
        self.assertIsNotNone(geom)
        self.assertEqual(geom["kind"], "straight")
        (x1, y1), (x2, y2) = geom["points"]
        self.assertLess(abs(y1 - 12), 2.0)
        self.assertLess(abs(y2 - 12), 2.0)
        self.assertLess(abs(x1 - 8), 4.0)
        self.assertLess(abs(x2 - 192), 4.0)
        self.assertLess(geom["width_px"], 6.5)
        self.assertIsNone(geom["dash"])
        self.assertFalse(geom["arrow_start"] or geom["arrow_end"])

    def test_arrowhead_detected_on_correct_end(self) -> None:
        crop = np.full((30, 200, 3), 255, np.uint8)
        cv2.arrowedLine(crop, (8, 15), (192, 15), (200, 60, 30), 3,
                        cv2.LINE_AA, tipLength=0.08)
        geom = _geometry_for(crop)
        self.assertIsNotNone(geom)
        # Points are normalised left→right, so the tip is arrow_end.
        self.assertTrue(geom["arrow_end"])
        self.assertFalse(geom["arrow_start"])

        crop2 = np.full((30, 200, 3), 255, np.uint8)
        cv2.arrowedLine(crop2, (192, 15), (8, 15), (200, 60, 30), 3,
                        cv2.LINE_AA, tipLength=0.08)
        geom2 = _geometry_for(crop2)
        self.assertIsNotNone(geom2)
        self.assertTrue(geom2["arrow_start"])
        self.assertFalse(geom2["arrow_end"])

    def test_dashed_line_detected(self) -> None:
        crop = np.full((24, 220, 3), 255, np.uint8)
        for x in range(10, 210, 22):
            cv2.line(crop, (x, 12), (x + 12, 12), (60, 60, 200), 3,
                     cv2.LINE_AA)
        geom = _geometry_for(crop)
        self.assertIsNotNone(geom)
        self.assertEqual(geom["dash"], "dash")

    def test_diagonal_line(self) -> None:
        crop = np.full((120, 160, 3), 255, np.uint8)
        cv2.line(crop, (10, 10), (150, 110), (30, 30, 30), 3, cv2.LINE_AA)
        geom = _geometry_for(crop)
        self.assertIsNotNone(geom)
        self.assertEqual(geom["kind"], "straight")

    def test_elbow_connector(self) -> None:
        crop = np.full((120, 160, 3), 255, np.uint8)
        cv2.line(crop, (10, 20), (120, 20), (90, 60, 30), 3)
        cv2.line(crop, (120, 20), (120, 110), (90, 60, 30), 3)
        geom = _geometry_for(crop)
        self.assertIsNotNone(geom)
        self.assertEqual(geom["kind"], "elbow")
        self.assertEqual(len(geom["points"]), 3)
        corner = geom["points"][1]
        self.assertLess(abs(corner[0] - 120), 4.0)
        self.assertLess(abs(corner[1] - 20), 4.0)

    def test_blob_and_curve_fall_back_to_raster(self) -> None:
        blob = np.full((80, 80, 3), 255, np.uint8)
        cv2.circle(blob, (40, 40), 25, (30, 30, 200), -1)
        self.assertIsNone(_geometry_for(blob))

        curve = np.full((120, 200, 3), 255, np.uint8)
        cv2.ellipse(curve, (100, 120), (90, 80), 0, 200, 340,
                    (30, 30, 30), 3, cv2.LINE_AA)
        self.assertIsNone(_geometry_for(curve))


class NativeLinePptxTests(unittest.TestCase):
    def test_straight_polyline_and_arrow_lines_build(self) -> None:
        import tempfile
        from build_pptx_from_layout import Builder

        layout = {
            "slide_size": {"width_in": 13.333, "height_in": 7.5},
            "source_width": 1280,
            "source_height": 720,
            "elements": [
                {"type": "line", "name": "l1",
                 "points": [100, 100, 500, 100],
                 "line": "#FF0000", "line_width": 2.0,
                 "arrow_end": True},
                {"type": "line", "name": "l2",
                 "points": [[100, 200], [400, 200], [400, 400]],
                 "box": [100, 200, 300, 200],
                 "line": "#0000FF", "line_width": 1.5, "dash": "dash"},
            ],
        }
        with tempfile.TemporaryDirectory() as td:
            out = Path(td) / "out.pptx"
            Builder(layout, out, Path(td)).build()
            self.assertTrue(out.exists())

            from pptx import Presentation
            prs = Presentation(str(out))
            shapes = list(prs.slides[0].shapes)
            names = {s.name for s in shapes}
            self.assertIn("l1", names)
            self.assertIn("l2", names)
            l1 = next(s for s in shapes if s.name == "l1")
            xml = l1.line._get_or_add_ln().xml
            self.assertIn("tailEnd", xml)


if __name__ == "__main__":
    unittest.main()
